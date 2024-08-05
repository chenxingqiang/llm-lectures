from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor

class PPTTemplateReader:
    def __init__(self, template_path):
        self.prs = Presentation(template_path)
        self.template_slides = self._categorize_slides()

    def _categorize_slides(self):
        categorized_slides = {'title': [], 'content': []}
        for slide in self.prs.slides:
            slide_info = self._get_slide_info(slide)
            if slide_info['title']:
                categorized_slides['title'].append(slide_info)
            if slide_info['content']:
                categorized_slides['content'].append(slide_info)
        return categorized_slides

    def _get_slide_info(self, slide):
        slide_info = {
            'title': self._get_shape_info(slide.shapes.title),
            'content': self._get_shape_info(self._find_placeholder(slide, 1))
        }
        return slide_info

    def _has_placeholder(self, slide, idx):
        try:
            slide.placeholders[idx]
            return True
        except KeyError:
            return False

    def _find_placeholder(self, slide, idx):
        """Find a placeholder by index and return it, or None if not found."""
        for placeholder in slide.placeholders:
            if placeholder.placeholder_format.idx == idx:
                return placeholder
        return None

    def _get_shape_info(self, shape):
        if not shape or not shape.has_text_frame:
            return None
        font_info = []
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                font_info.append({
                    'font_size': run.font.size,
                    'font_color': self._get_font_color(run.font),
                    'font_bold': run.font.bold,
                    'font_italic': run.font.italic,
                    'font_underline': run.font.underline,
                    'font_name': run.font.name
                })
        return font_info

    def _get_font_color(self, font):
        if font.color is None:
            return RGBColor(0, 0, 0)  # Default to black if no color specified
        if font.color.type == 1:  # RGB color
            return font.color.rgb
        elif font.color.type == 2:  # Theme color
            return RGBColor(0, 0, 0)  # Default to black
        else:
            return RGBColor(0, 0, 0)  # Default to black

    def create_new_slide(self, category, title_text, content_text):
        if category not in self.template_slides:
            raise ValueError(f"No slides found for category '{category}'")

        slide_info_list = self.template_slides[category]
        if not slide_info_list:
            raise ValueError(f"No slides available for category '{category}'")

        # Select a template slide based on the category
        template_slide_info = slide_info_list[0]  # You can add logic to choose a specific slide

        slide_layout = self.prs.slide_layouts[1]  # Assuming layout 1 has the necessary placeholders
        slide = self.prs.slides.add_slide(slide_layout)
        
        title_info = template_slide_info['title']
        content_info = template_slide_info['content']
        
        title = slide.shapes.title
        title.text = title_text
        self._apply_font_info(title.text_frame, title_info)
        
        content_placeholder = self._find_placeholder(slide, 1)
        if content_placeholder:
            content = content_placeholder
            content.text = content_text
            self._apply_font_info(content.text_frame, content_info)
        
        return slide

    def _apply_font_info(self, text_frame, font_info):
        for i, paragraph in enumerate(text_frame.paragraphs):
            if i < len(font_info):
                for run in paragraph.runs:
                    run.font.size = font_info[i]['font_size']
                    run.font.color.rgb = font_info[i]['font_color']
                    run.font.bold = font_info[i]['font_bold']
                    run.font.italic = font_info[i]['font_italic']
                    run.font.underline = font_info[i]['font_underline']
                    run.font.name = font_info[i]['font_name']

    def save_presentation(self, output_path):
        self.prs.save(output_path)

# 使用示例
template_path = 'template.pptx'
output_path = 'updated_presentation.pptx'
ppt_reader = PPTTemplateReader(template_path)

# 定义要插入的内容及其类别
slides_content = [
    {
        "category": "title",
        "title": "Introduction",
        "content": "This is the introduction slide."
    },
    {
        "category": "content",
        "title": "Methodology",
        "content": "This slide explains the methodology."
    },
    {
        "category": "content",
        "title": "Results",
        "content": "Here are the results of our analysis."
    },
    {
        "category": "title",
        "title": "Conclusion",
        "content": "This is the conclusion slide."
    }
]

# 创建新的幻灯片并插入内容
for slide_content in slides_content:
    ppt_reader.create_new_slide(slide_content["category"], slide_content["title"], slide_content["content"])

# 保存修改后的PPT
ppt_reader.save_presentation(output_path)

print(f"PPT has been saved as {output_path}")