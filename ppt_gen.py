import markdown
from bs4 import BeautifulSoup
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

import re
import os

class MarkdownReader:
    def __init__(self, markdown_path):
        self.markdown_path = markdown_path
        with open(markdown_path, "r", encoding="utf-8") as file:
            self.md_content = file.read()
        self.html_content = markdown.markdown(self.md_content)
        self.soup = BeautifulSoup(self.html_content, "html.parser")

    def get_title(self):
        title = self.soup.find("h1")
        return title.text if title else ""

    def get_subtitle(self):
        subtitle = self.soup.find("h2")
        return subtitle.text if subtitle else ""

    def get_sections(self):
        sections = []
        current_section = None
        for elem in self.soup.find_all(["h3", "p", "ul", "ol"]):
            if elem.name == "h3":
                if current_section:
                    sections.append(current_section)
                current_section = {"title": elem.text, "content": []}
            elif current_section:
                if elem.name == "p":
                    current_section["content"].append(elem.text)
                elif elem.name in ["ul", "ol"]:
                    list_items = [li.text for li in elem.find_all("li")]
                    current_section["content"].append("\n".join(f"• {item}" for item in list_items))
        if current_section:
            sections.append(current_section)
        return sections

class PPTTemplateReader:
    def __init__(self, template_path):
        self.prs = Presentation(template_path)
        self.template_slides = self._categorize_slides()

    def _categorize_slides(self):
        categorized_slides = {"title": [], "toc": [], "content": [], "section": [], "closing": []}
        for slide_layout in self.prs.slide_master.slide_layouts:
            slide_layout_name = slide_layout.name.lower()
            print(f"Slide layout name: {slide_layout_name}")  # Debug information
            if "title" in slide_layout_name:
                categorized_slides["title"].append(slide_layout)
            elif "toc" in slide_layout_name or "table of contents" in slide_layout_name:
                categorized_slides["toc"].append(slide_layout)
            elif "content" in slide_layout_name:
                categorized_slides["content"].append(slide_layout)
            elif "section" in slide_layout_name:
                categorized_slides["section"].append(slide_layout)
            elif "closing" in slide_layout_name:
                categorized_slides["closing"].append(slide_layout)
        print(f"Categorized slides: {categorized_slides}")  # Debug information
        return categorized_slides

    def create_new_slide(self, category, title_text, content_text):
        if category not in self.template_slides:
            raise ValueError(f"No slides found for category '{category}'")

        slide_layout_list = self.template_slides[category]
        if not slide_layout_list:
            raise ValueError(f"No slides available for category '{category}'")

        # Select a template slide based on the category
        slide_layout = slide_layout_list[
            0
        ]  # You can add logic to choose a specific slide

        slide = self.prs.slides.add_slide(slide_layout)

        title_placeholder = slide.shapes.title
        if title_placeholder:
            title_placeholder.text = title_text

        content_placeholder = self._find_placeholder(slide, 1)
        if content_placeholder:
            self._set_content(content_placeholder, content_text)

        return slide

    def _find_placeholder(self, slide, idx):
        for placeholder in slide.placeholders:
            if placeholder.placeholder_format.idx == idx:
                return placeholder
        return None

    def _set_content(self, placeholder, content_text):
        text_frame = placeholder.text_frame
        text_frame.clear()

        for line in content_text.split("\n"):
            p = text_frame.add_paragraph()
            p.text = line
            p.space_before = Pt(3)
            p.space_after = Pt(5)
            p.font.size = Pt(42)

        # Adjust the font size to fit text within the placeholder
        text_frame.auto_size = True
        current_font_size = Pt(42)

        while text_frame.text != content_text:
            current_font_size = Pt(current_font_size.pt - 1)
            for paragraph in text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = current_font_size

            if current_font_size.pt <= 36:  # Minimum font size limit
                break


    def save_presentation(self, output_path):
        self.prs.save(output_path)

def remove_non_alphanumeric(s):
    # 只保留中文、英文和数字字符
    filtered_string = re.sub(r'[^\u4e00-\u9fa5a-zA-Z0-9]', '', s)
    return filtered_string

class DetailedPPTGenerator:
    def __init__(self, template_reader, markdown_reader, output_dir):
        self.template_reader = template_reader
        self.markdown_reader = markdown_reader
        self.output_dir = output_dir
        self.title = markdown_reader.get_title()
        self.subtitle = markdown_reader.get_subtitle()

    def _add_slide(self, category, title, content):
        self.template_reader.create_new_slide(category, title, content)

    def generate_ppt(self):
        # Title Slide
        self._add_slide("title", self.title, self.subtitle)

        # Table of Contents
        sections = self.markdown_reader.get_sections()
        toc_content = "\n".join(
            [
                f"{i+1}. {section['title']}"
                for i, section in enumerate(sections)
            ]
        )
        self._add_slide("toc", "目录", toc_content)

        # Section Name Slide for each section
        for section in sections:
            self._add_slide("section", section["title"], "")

            # Content Slides within the section
            self._add_slide("content", section["title"], "\n".join(section["content"]))

        # Ensure output directory exists
        if not os.path.exists(self.output_dir):
            os.makedirs(self.output_dir)

        # Save the presentation
        markdown_filename = os.path.basename(self.markdown_reader.markdown_path)
        markdown_name = os.path.splitext(markdown_filename)[0]
        sanitized_title = remove_non_alphanumeric(self.title)
        sanitized_markdown_name = remove_non_alphanumeric(markdown_name)
        filename = f"{sanitized_markdown_name}-{sanitized_title}.pptx"
        output_path = os.path.join(self.output_dir, filename)
        self.template_reader.save_presentation(output_path)

